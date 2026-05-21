"""
html_to_pptx.py
Converts itweb_keynote_v5.html → itweb_keynote_v5.pptx

Workflow:
  1. Playwright screenshots every slide at 1920×1080 (UI overlays hidden).
  2. python-pptx builds the PPTX:
     - Slide 1  → template cover slide (slide index 0)
     - Slides 2-N → template content slide (slide index 1) duplicated
  3. Each screenshot fills the full slide.
  4. Banner shapes from the slide layout are copied into each slide at the
     top of the z-order so they appear on top of the screenshot.
"""

import copy, io, pathlib, sys, time
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from playwright.sync_api import sync_playwright

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE      = pathlib.Path('C:/Users/steph/OneDrive/Documents/GitHub/ITWeb/analysis')
HTML_FILE = BASE / 'itweb_keynote_v5.html'
TEMPLATE  = pathlib.Path('C:/Users/steph/AppData/Local/Temp/template.pptx')
OUTPUT    = BASE / 'itweb_keynote_v5_from_html.pptx'
OUTPUT_TMP = pathlib.Path('C:/Users/steph/AppData/Local/Temp/itweb_keynote_v5.pptx')

VIEWPORT_W   = 1920
VIEWPORT_H   = 1080
SLIDE_COUNT  = 20          # confirmed from HTML inspection
TRANS_WAIT   = 500         # ms to wait after activating a slide

# ── Step 1: Screenshot every slide ───────────────────────────────────────────
print('Step 1: Screenshotting HTML slides ...')
screenshots = []   # list of PNG bytes, one per slide

with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    page = browser.new_page(viewport={'width': VIEWPORT_W, 'height': VIEWPORT_H})
    page.goto(HTML_FILE.as_uri())
    page.wait_for_load_state('networkidle', timeout=30_000)
    time.sleep(1)  # let animations settle

    # Hide all chrome/UI overlays so only the slide content is visible
    page.evaluate("""() => {
        const hide = [
            '.nav-controls', '.top-bar', '#notes-toggle', '#notes-panel',
            '.kb-hint', '.fs-btn', '.mouse-spotlight', '#notes-toggle',
            '.slide-counter'
        ];
        hide.forEach(sel => {
            document.querySelectorAll(sel).forEach(el => el.style.display = 'none');
        });
    }""")

    for i in range(SLIDE_COUNT):
        # Activate slide i via the JS function defined in the HTML
        page.evaluate(f'goToSlide({i + 1})')
        page.wait_for_timeout(TRANS_WAIT)
        png = page.screenshot(full_page=False, type='png')
        screenshots.append(png)
        print(f'  Captured slide {i + 1}/{SLIDE_COUNT}  ({len(png):,} bytes)')

    browser.close()

print(f'  {len(screenshots)} screenshots captured.\n')

# ── Step 2: Build PPTX ───────────────────────────────────────────────────────
print('Step 2: Building PPTX ...')

tpl = Presentation(str(TEMPLATE))
slide_w = tpl.slide_width    # 9144000 emu
slide_h = tpl.slide_height   # 5143500 emu

# The template has exactly 2 slides:
#   index 0 → cover layout  (top + bottom banner)
#   index 1 → content layout (bottom banner only)
tpl_cover_slide   = tpl.slides[0]
tpl_content_slide = tpl.slides[1]


def _spTree(slide):
    return slide.shapes._spTree


def add_screenshot_to_slide(slide, png_bytes):
    """Insert screenshot as a full-slide background image at z-order 0."""
    img = io.BytesIO(png_bytes)
    pic = slide.shapes.add_picture(img, 0, 0, slide_w, slide_h)
    # Move to back: position right after the two mandatory group-property elements
    tree = _spTree(slide)
    tree.remove(pic._element)
    # Insert after nvGrpSpPr (idx 0) and grpSpPr (idx 1)
    tree.insert(2, pic._element)


def copy_banner_shapes(slide):
    """Copy PICTURE shapes from the slide layout into the slide at top z-order."""
    layout = slide.slide_layout
    tree = _spTree(slide)
    for shape in layout.shapes:
        if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE == 13
            elem = copy.deepcopy(shape._element)
            tree.append(elem)        # append = highest z-order (drawn last = on top)


def duplicate_content_slide(prs, template_slide):
    """Create a new slide that is a copy of template_slide."""
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # Clear auto-populated placeholders from the new slide
    new_tree = _spTree(new_slide)
    for child in list(new_tree):
        tag = child.tag
        if tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
            new_tree.remove(child)
    # Copy all non-property children from template (none expected, but just in case)
    tpl_tree = _spTree(template_slide)
    for child in tpl_tree:
        if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
            new_tree.append(copy.deepcopy(child))
    return new_slide


def remove_slide(prs, index):
    """Remove slide at index from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn('r:id'))
    sldIdLst.remove(sldId)
    # Remove the relationship and drop the slide part
    prs.part.drop_rel(rId)


# ── Slide 1: cover ────────────────────────────────────────────────────────────
print('  Processing slide 1 (cover) ...')
cover = tpl_cover_slide
add_screenshot_to_slide(cover, screenshots[0])
copy_banner_shapes(cover)

# ── Slides 2–20: content (duplicate template slide 2) ─────────────────────────
for i in range(1, SLIDE_COUNT):
    print(f'  Processing slide {i + 1}/{SLIDE_COUNT} ...')
    new_slide = duplicate_content_slide(tpl, tpl_content_slide)
    add_screenshot_to_slide(new_slide, screenshots[i])
    copy_banner_shapes(new_slide)

# ── Remove the original blank template content slide ─────────────────────────
# It is now at index 1 (slide 2), pushed there; our new slides were appended
# after it. Actually add_slide appends, so order is:
#   index 0: cover (modified in place)  ← keep
#   index 1: original blank content     ← remove
#   index 2-20: our 19 new content slides ← keep
print('  Removing blank template content slide ...')
remove_slide(tpl, 1)

# ── Save ──────────────────────────────────────────────────────────────────────
import shutil
tpl.save(str(OUTPUT_TMP))
shutil.copy2(str(OUTPUT_TMP), str(OUTPUT))
print(f'\nDone: {OUTPUT}')
print(f'Slides in output: {len(tpl.slides)}')
